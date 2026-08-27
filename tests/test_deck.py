from pathlib import Path

from pptx import Presentation

from deck.build_deck import parse_packet

DECK_DIR = Path(__file__).resolve().parents[1] / "deck"
PACKET = DECK_DIR / "DECK_PACKET.md"


def test_packet_has_30_plus_slides():
    slides = parse_packet(PACKET)
    assert len(slides) >= 30


def test_every_slide_has_title():
    for s in parse_packet(PACKET):
        assert s["title"], f"slide {s['num']} missing title"


def test_every_slide_has_notes():
    slides = parse_packet(PACKET)
    assert all(s["notes"] for s in slides), "at least one slide has empty narration"


def test_slide_numbers_contiguous_from_one():
    nums = [s["num"] for s in parse_packet(PACKET)]
    assert nums == list(range(1, len(nums) + 1))


def test_all_figures_exist_in_deck_assets():
    assets = DECK_DIR / "deck_assets"
    missing = []
    for s in parse_packet(PACKET):
        if s["figure"]:
            path = assets / s["figure"]["path"]
            if not path.exists():
                missing.append(str(path))
    assert not missing, f"missing figures: {missing}"


def test_rendered_deck_opens_as_16x9_with_notes(tmp_path):
    from deck.build_deck import render

    out = tmp_path / "deck.pptx"
    slides = parse_packet(PACKET)
    render(DECK_DIR, out, slides)
    prs = Presentation(str(out))
    assert len(prs.slides) >= 30
    assert prs.slide_width > prs.slide_height
    with_notes = sum(
        1
        for s in prs.slides
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
    )
    assert with_notes == len(prs.slides)
